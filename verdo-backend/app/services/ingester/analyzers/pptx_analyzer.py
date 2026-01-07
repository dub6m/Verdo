from pptx import Presentation

from app.services.ingester.core.shape_classifier import ShapeClassifier

# --- Classes ------------------------------------------------------------

class PptxAnalyzer:
	def __init__(self):
		self.shapeClassifier = ShapeClassifier()

	def analyze(self, filePath: str):
		prs = Presentation(filePath)
		results = []
		for slideNum, slide in enumerate(prs.slides):
			elements = []
			for shape in slide.shapes:
				elemType = self.shapeClassifier.classifyShape(shape)
				element = {
					'slide_number': slideNum + 1,
					'raw_type': shape.shape_type,
					'classified_type': elemType,
					'position': (shape.top, shape.left),
					'size': (shape.width, shape.height),
					'content': shape.text[:100] if elemType == 'text' and shape.has_text_frame else f"<{elemType.upper()}>"
				}
				elements.append(element)
			elements.sort(key=lambda x: (x['position'][0], x['position'][1]))
			results.append({'slide_number': slideNum + 1, 'elements': elements})
		return results

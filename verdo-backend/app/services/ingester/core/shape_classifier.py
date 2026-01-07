from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- Classes ------------------------------------------------------------

class ShapeClassifier:
	# Classify PowerPoint shapes into text, image, table, etc.
	@staticmethod
	def classifyShape(shape):
		t = shape.shape_type
		if t == MSO_SHAPE_TYPE.PICTURE: return 'image'
		if t == MSO_SHAPE_TYPE.TABLE: return 'table'
		if t == MSO_SHAPE_TYPE.CHART: return 'chart'
		# Math can be either an embedded OLE object (MathType/legacy)
		# or inline OMML within a text box/run. Detect both.
		if t == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: return 'equation'
		try:
			if ShapeClassifier._hasOmml(shape):
				return 'equation'
		except Exception:
			pass
		if t == 15: return 'diagram'
		if t == MSO_SHAPE_TYPE.GROUP:
			# Groups can wrap OLE/OMML children; treat as equation if present
			try:
				if ShapeClassifier._hasOmml(shape) or ShapeClassifier._hasOle(shape):
					return 'equation'
			except Exception:
				pass
			return 'group'
		if t == 17: return 'text'
		if t == MSO_SHAPE_TYPE.PLACEHOLDER:
			try:
				_ = shape.image; return 'image'
			except AttributeError: pass
			# If placeholder contains OMML, treat it as equation
			try:
				if ShapeClassifier._hasOmml(shape):
					return 'equation'
			except Exception:
				pass
			if shape.has_text_frame and shape.text.strip(): return 'text'
			return 'unknown'
		# Text fallback (but prefer equation if OMML is present)
		if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
			try:
				if ShapeClassifier._hasOmml(shape):
					return 'equation'
			except Exception:
				pass
			return 'text'
		return 'unknown'

	@staticmethod
	def _hasOmml(shape) -> bool:
		# Lightweight check for OMML in a shape.
		# Looks for oMath/oMathPara elements via xpath if available,
		# otherwise searches the element XML string.
		el = getattr(shape, '_element', None)
		if el is None:
			el = getattr(shape, 'element', None)
		if el is None:
			return False
		# Preferred: xpath lookup
		if hasattr(el, 'xpath'):
			try:
				ns = {"m": "http://schemas.openxmlformats.org/officeDocument/2006/math"}
				nodes = el.xpath('.//m:oMath | .//m:oMathPara', namespaces=ns)
				if nodes:
					return True
			except Exception:
				pass
		# Fallback: string search on XML
		xml = getattr(el, 'xml', None)
		if isinstance(xml, str) and ('oMath' in xml or 'oMathPara' in xml):
			return True
		# Last resort: scan text runs for embedded OMML fragments
		tf = getattr(shape, 'text_frame', None)
		if tf is not None:
			try:
				for p in tf.paragraphs:
					for r in p.runs:
						rx = getattr(r, '_r', None)
						rxml = getattr(rx, 'xml', None)
						if isinstance(rxml, str) and ('oMath' in rxml or 'oMathPara' in rxml):
							return True
			except Exception:
				pass
		return False

	@staticmethod
	def _hasOle(shape) -> bool:
		# Detect embedded OLE objects within a shape (including inside groups).
		el = getattr(shape, '_element', None)
		if el is None:
			el = getattr(shape, 'element', None)
		if el is None:
			return False
		if hasattr(el, 'xpath'):
			try:
				ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
				nodes = el.xpath('.//p:oleObj', namespaces=ns)
				if nodes:
					return True
			except Exception:
				pass
		xml = getattr(el, 'xml', None)
		if isinstance(xml, str) and 'oleObj' in xml:
			return True
		return False

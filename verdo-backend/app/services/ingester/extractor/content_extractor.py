import fitz  # PyMuPDF
from pptx import Presentation

from app.services.ingester.handlers.default_handler import DefaultHandler
from app.services.ingester.handlers.formula_handler import FormulaHandler
from app.services.ingester.handlers.image_handler import ImageHandler
from app.services.ingester.handlers.table_handler import TableHandler
from app.services.ingester.handlers.text_handler import TextHandler

# --- Classes ------------------------------------------------------------

class ContentExtractor:
	# Routes analyzed elements to their correct handlers.

	def __init__(self, verbose=True, apiKey=None, llmClient=None, enableCache=True):
		# Verbose: Whether to print processing messages
		self.verbose = verbose
		self.llm = llmClient
		self.textHandler = TextHandler()
		
		# Initialize dependencies first
		self.tableHandler = TableHandler(
			apiKey=apiKey,
			llmClient=llmClient,
			enableGptFallback=True,
			verbose=verbose
		)
		self.formulaHandler = FormulaHandler(
			apiKey=apiKey,
			llmClient=llmClient,
			verbose=verbose
		)
		
		# Helper for Image Handler (avoid circular dep by passing instances)
		self.imageHandler = ImageHandler(
			apiKey=apiKey,
			llmClient=llmClient,
			enableCache=enableCache,
			verbose=verbose,
			tableHandler=self.tableHandler,
			formulaHandler=self.formulaHandler
		)

		self.defaultHandler = DefaultHandler()

	# ===== Universal Entry Point =====

	# Extract content from analyzed elements
	def extract(self, filePath, analyzedElements):
		if filePath.endswith(".pptx"):
			return self.extractFromPptx(filePath, analyzedElements)
		elif filePath.endswith(".pdf"):
			return self.extractFromPdf(filePath, analyzedElements)
		else:
			raise ValueError(f"Unsupported file type: {filePath}")

	# ===== PPTX Extraction =====

	# Extract content from PowerPoint slides with parallel image processing
	def extractFromPptx(self, filePath, analyzedElements):
		prs = Presentation(filePath)

		# First pass: submit all image processing tasks and collect futures
		imageFutures = []  # List of (slideNum, elemIndex, future)
		elementsMetadata = []  # List of (slideNum, elemList)

		for slideData in analyzedElements:
			slideNum = slideData['slide_number']
			slide = prs.slides[slideNum - 1]
			slideElements = []

			for elemIndex, elem in enumerate(slideData['elements']):
				elemType = elem['classified_type']
				shape = self._findShape(slide, elem['position'])

				if shape is None:
					if self.verbose:
						print(f"⚠️  No shape found at {elem['position']}")
					continue

				# Store element metadata
				elemData = {
					'type': elemType,
					'position': elem['position'],
					'size': elem['size'],
					'content': None,  # Will be filled later
					'shape': shape,
					'handler': self._getHandler(elemType)
				}
				slideElements.append(elemData)

				# For images, submit async task immediately
				if elemType == 'image' and hasattr(self.imageHandler, 'handlePptxAsync'):
					try:
						future = self.imageHandler.handlePptxAsync(shape)
						imageFutures.append((slideNum, len(slideElements) - 1, future))
					except Exception as e:
						if self.verbose:
							print(f"⚠️  Error submitting image task: {e}")

			elementsMetadata.append((slideNum, slideElements))

		# Second pass: process non-image elements
		for slideNum, slideElements in elementsMetadata:
			for elemData in slideElements:
				if elemData['type'] != 'image':
					handler = elemData['handler']
					try:
						# Pass filePath and slideNum so slide-level OMML fallback can work
						content = handler.handlePptx(
							elemData['shape'],
							filePath=filePath,
							slideNum=slideNum,
						)
					except AttributeError:
						# Handler doesn't implement handlePptx (e.g., DefaultHandler)
						content = handler.handle()
					elemData['content'] = content

		# Third pass: wait for all image futures to complete
		for slideNum, elemIndex, future in imageFutures:
			# Find the corresponding element
			for sNum, slideElements in elementsMetadata:
				if sNum == slideNum:
					slideElements[elemIndex]['content'] = future.result()
					break

		# Build final results
		results = []
		for slideNum, slideElements in elementsMetadata:
			extracted = []
			for elemData in slideElements:
				extracted.append({
					'type': elemData['type'],
					'position': elemData['position'],
					'size': elemData['size'],
					'content': elemData['content']
				})
			results.append({'slide_number': slideNum, 'elements': extracted})

		return results

	# ===== PDF Extraction =====

	# Extract content from PDF pages with parallel image processing
	def extractFromPdf(self, filePath, analyzedElements, dpi=250):
		doc = fitz.open(filePath)

		# First pass: submit all image processing tasks and collect futures
		imageFutures = []  # List of (pageNum, detIndex, future)
		elementsMetadata = []  # List of (pageNum, detList)

		for pageData in analyzedElements:
			pageNum = pageData['page_number']
			page = doc[pageNum - 1]

			# Calculate scale factor between YOLO and PDF coordinates
			detectionW = int(page.rect.width * dpi / 72)
			detectionH = int(page.rect.height * dpi / 72)
			scale = ((page.rect.width / detectionW) + (page.rect.height / detectionH)) / 2

			pageElements = []
			for detIndex, det in enumerate(pageData['detections']):
				cls, bbox, conf = det['class_name'], det['bbox'], det['conf']

				# Store element metadata
				elemData = {
					'type': cls,
					'bbox': bbox,
					'confidence': conf,
					'content': None,  # Will be filled later
					'handler': self._getHandler(cls),
					'page': page,
					'scale': scale
				}
				pageElements.append(elemData)

				# For images/figures, submit async task immediately
				if cls in ['image', 'figure'] and hasattr(self.imageHandler, 'handlePdfAsync'):
					try:
						future = self.imageHandler.handlePdfAsync(page, bbox, scale)
						imageFutures.append((pageNum, len(pageElements) - 1, future))
					except Exception as e:
						if self.verbose:
							print(f"⚠️  Error submitting image task: {e}")

			elementsMetadata.append((pageNum, pageElements))

		# Second pass: process non-image elements
		for pageNum, pageElements in elementsMetadata:
			for elemData in pageElements:
				if elemData['type'] not in ['image', 'figure']:
					handler = elemData['handler']
					try:
						content = handler.handlePdf(elemData['page'], elemData['bbox'], elemData['scale'])
					except AttributeError:
						# Handler doesn't implement handlePdf (e.g., DefaultHandler)
						content = handler.handle()
					elemData['content'] = content

		# Third pass: wait for all image futures to complete
		for pageNum, detIndex, future in imageFutures:
			# Find the corresponding element
			for pNum, pageElements in elementsMetadata:
				if pNum == pageNum:
					pageElements[detIndex]['content'] = future.result()
					break

		# Build final results
		results = []
		for pageNum, pageElements in elementsMetadata:
			extracted = []
			for elemData in pageElements:
				extracted.append({
					'type': elemData['type'],
					'bbox': elemData['bbox'],
					'confidence': elemData['confidence'],
					'content': elemData['content']
				})
			results.append({'page_number': pageNum, 'elements': extracted})

		doc.close()
		return results

	# ===== Utilities =====

	# Find shape in slide by position
	def _findShape(self, slide, position):
		top, left = position
		for shape in slide.shapes:
			if shape.top == top and shape.left == left:
				return shape
		return None

	# Route element type to appropriate handler
	def _getHandler(self, elemType):
		handlerMap = {
			# Text types
			'text': self.textHandler,
			'plain text': self.textHandler,
			'title': self.textHandler,

			# Image types
			'image': self.imageHandler,
			'figure': self.imageHandler,

			# Other types
			'table': self.tableHandler,
			'equation': self.formulaHandler,
			'isolate_formula': self.formulaHandler,
			# 'chart': self.chartHandler,
			# 'diagram': self.diagramHandler,
			# 'group': self.groupHandler,
		}

		return handlerMap.get(elemType, self.defaultHandler)
